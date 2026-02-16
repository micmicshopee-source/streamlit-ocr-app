# -*- coding: utf-8 -*-
"""
PDF 萬能轉換工具模組
支援：PDF → Excel, PPT, 圖片 (JPG/PNG), Word
      圖片 / Word / Excel / PPT → PDF
含 AI OCR 模式：掃描檔 PDF 轉 Word（使用 Gemini Vision）
"""

from __future__ import annotations

import base64
import io
import json
import os
import re
import subprocess
import tempfile
import zipfile
from typing import List, Optional, Tuple, Union

# 可選依賴：按需導入
_pdfplumber = None
_pymupdf = None
_pptx = None
_pdf2image = None
_pdf2docx = None
_pandas = None
_openpyxl = None
_pil = None
_python_docx = None
_requests = None


def _pil_lanczos():
    """
    取得與 Pillow 版本相容的縮圖 resample 方式。
    新版 Pillow 有 Image.Resampling.LANCZOS，舊版只有 Image.LANCZOS / ANTIALIAS。
    """
    global _pil
    if not _pil:
        return None
    # Pillow >= 9.1
    if hasattr(_pil, "Resampling"):
        try:
            return _pil.Resampling.LANCZOS  # type: ignore[attr-defined]
        except Exception:
            pass
    # 舊版 Pillow
    if hasattr(_pil, "LANCZOS"):
        return _pil.LANCZOS  # type: ignore[attr-defined]
    if hasattr(_pil, "ANTIALIAS"):
        return _pil.ANTIALIAS  # type: ignore[attr-defined]
    return None


def _safe_imports():
    """延遲導入，避免 import 時即失敗。"""
    global _pdfplumber, _pptx, _pdf2image, _pdf2docx, _pandas, _openpyxl, _pil, _python_docx, _requests, _pymupdf
    if _pdfplumber is None:
        try:
            import pdfplumber
            _pdfplumber = pdfplumber
        except ImportError:
            _pdfplumber = False
    if _pptx is None:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            _pptx = (Presentation, Inches, Pt)
        except ImportError:
            _pptx = False
    if _pdf2image is None:
        try:
            from pdf2image import convert_from_bytes
            _pdf2image = convert_from_bytes
        except ImportError:
            _pdf2image = False
    if _pdf2docx is None:
        try:
            from pdf2docx import Converter
            _pdf2docx = Converter
        except ImportError:
            _pdf2docx = False
    if _pandas is None:
        try:
            import pandas as pd
            _pandas = pd
        except ImportError:
            _pandas = False
    if _openpyxl is None:
        try:
            from openpyxl import load_workbook
            from openpyxl.styles import Alignment, Font, Border, Side
            _openpyxl = (load_workbook, Alignment, Font, Border, Side)
        except ImportError:
            _openpyxl = False
    if _pil is None:
        try:
            from PIL import Image
            _pil = Image
        except ImportError:
            _pil = False
    if _python_docx is None:
        try:
            from docx import Document
            from docx.shared import Pt
            _python_docx = (Document, Pt)
        except ImportError:
            _python_docx = False
    if _requests is None:
        try:
            import requests
            _requests = requests
        except ImportError:
            _requests = False
    if _pymupdf is None:
        try:
            import fitz
            _pymupdf = fitz
        except ImportError:
            _pymupdf = False


def _office_to_pdf_via_libreoffice(
    file_bytes: bytes,
    ext: str,
    progress_callback=None,
) -> Tuple[Optional[bytes], Optional[str]]:
    """
    使用 LibreOffice 將 Office 檔轉為 PDF。
    需系統安裝 LibreOffice：sudo apt install libreoffice
    Returns: (pdf_bytes, error_message)
    """
    _cmd = None
    for cmd in ("libreoffice", "soffice"):
        try:
            subprocess.run([cmd, "--version"], capture_output=True, timeout=5)
            _cmd = cmd
            break
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    if not _cmd:
        return None, "未安裝 LibreOffice，請執行：sudo apt install libreoffice"

    try:
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as f:
            f.write(file_bytes)
            input_path = f.name
        out_dir = tempfile.mkdtemp()
        try:
            proc = subprocess.run(
                [_cmd, "--headless", "--convert-to", "pdf", "--outdir", out_dir, input_path],
                capture_output=True,
                text=True,
                timeout=120,
            )
            if proc.returncode != 0:
                return None, f"LibreOffice 轉換失敗：{proc.stderr or proc.stdout or '未知錯誤'}"
            base_name = os.path.splitext(os.path.basename(input_path))[0]
            pdf_path = os.path.join(out_dir, f"{base_name}.pdf")
            if not os.path.isfile(pdf_path):
                return None, "轉換後未產生 PDF 檔案"
            with open(pdf_path, "rb") as f:
                return f.read(), None
        finally:
            try:
                for f in os.listdir(out_dir):
                    os.remove(os.path.join(out_dir, f))
                os.rmdir(out_dir)
            except Exception:
                pass
            try:
                os.remove(input_path)
            except Exception:
                pass
    except subprocess.TimeoutExpired:
        return None, "轉換逾時（超過 120 秒）"
    except Exception as e:
        return None, f"轉換失敗：{str(e)}"


def word_to_pdf(docx_bytes: bytes, progress_callback=None) -> Tuple[Optional[bytes], Optional[str]]:
    """Word (.docx) 轉 PDF。需 LibreOffice。"""
    return _office_to_pdf_via_libreoffice(docx_bytes, "docx", progress_callback)


def excel_to_pdf(xlsx_bytes: bytes, progress_callback=None) -> Tuple[Optional[bytes], Optional[str]]:
    """Excel (.xlsx) 轉 PDF。需 LibreOffice。"""
    return _office_to_pdf_via_libreoffice(xlsx_bytes, "xlsx", progress_callback)


def ppt_to_pdf(pptx_bytes: bytes, progress_callback=None) -> Tuple[Optional[bytes], Optional[str]]:
    """PPT (.pptx) 轉 PDF。需 LibreOffice。"""
    return _office_to_pdf_via_libreoffice(pptx_bytes, "pptx", progress_callback)


def images_to_pdf(
    image_bytes_list: List[bytes],
    progress_callback=None,
) -> Tuple[Optional[bytes], Optional[str]]:
    """
    將多張圖片（JPG/PNG）合併為單一 PDF。
    圖片依上傳順序排列，每張一頁。
    Returns: (pdf_bytes, error_message)
    """
    _safe_imports()
    if not _pymupdf:
        return None, "未安裝 pymupdf，請執行：pip install pymupdf"
    if not _pil:
        return None, "未安裝 Pillow"

    if not image_bytes_list:
        return None, "請至少上傳一張圖片"

    try:
        fitz = _pymupdf
        doc = fitz.open()
        total = len(image_bytes_list)

        for i, img_bytes in enumerate(image_bytes_list):
            if progress_callback:
                progress_callback((i + 1) / total)

            # 用 PIL 偵測格式（pymupdf 需 png/jpeg）
            img = _pil.open(io.BytesIO(img_bytes))
            fmt = (img.format or "PNG").upper()
            ext = "png" if fmt in ("PNG", "PPM", "GIF", "BMP") else "jpeg"

            img_doc = fitz.open(stream=img_bytes, filetype=ext)
            pdf_bytes = img_doc.convert_to_pdf()
            img_doc.close()

            img_pdf = fitz.open(stream=pdf_bytes, filetype="pdf")
            doc.insert_pdf(img_pdf)
            img_pdf.close()

        result = doc.tobytes(garbage=4, deflate=True)
        doc.close()
        return result, None
    except Exception as e:
        err_msg = str(e)
        return None, f"轉換失敗：{err_msg}"


def pdf_to_excel(pdf_bytes: bytes, progress_callback=None) -> Tuple[Optional[bytes], Optional[str]]:
    """
    使用 pdfplumber 提取 PDF 表格，導出為 .xlsx。
    多個表格分別放在不同 Sheet。
    Returns: (xlsx_bytes, error_message)
    """
    _safe_imports()
    if not _pdfplumber:
        return None, "未安裝 pdfplumber，請執行：pip install pdfplumber"
    if not _pandas:
        return None, "未安裝 pandas"
    if not _openpyxl:
        return None, "未安裝 openpyxl"

    try:
        import pandas as pd
        with io.BytesIO(pdf_bytes) as f:
            pdf = _pdfplumber.open(f)
            tables_by_page = []
            total_pages = len(pdf.pages)
            for i, page in enumerate(pdf.pages):
                if progress_callback:
                    progress_callback((i + 1) / total_pages)
                tables = page.extract_tables()
                if tables:
                    for t_idx, tbl in enumerate(tables):
                        if tbl and len(tbl) > 0:
                            df = pd.DataFrame(tbl[1:], columns=tbl[0] if tbl[0] else None)
                            if df.empty and tbl:
                                df = pd.DataFrame(tbl)
                            tables_by_page.append((f"Page{i+1}_Table{t_idx+1}", df))
            pdf.close()

        if not tables_by_page:
            return None, "PDF 中未偵測到表格數據"

        buf = io.BytesIO()
        with pd.ExcelWriter(buf, engine="openpyxl") as writer:
            for sheet_name, df in tables_by_page:
                safe_name = sheet_name[:31]  # Excel sheet 名稱限制 31 字元
                df.to_excel(writer, sheet_name=safe_name, index=False)
        buf.seek(0)
        return buf.read(), None
    except Exception as e:
        err_msg = str(e)
        if "encrypted" in err_msg.lower() or "password" in err_msg.lower():
            return None, "PDF 已加密或受密碼保護，無法讀取"
        if "corrupt" in err_msg.lower() or "invalid" in err_msg.lower():
            return None, "PDF 格式損毀或無效"
        return None, f"轉換失敗：{err_msg}"


def pdf_to_ppt(pdf_bytes: bytes, progress_callback=None) -> Tuple[Optional[bytes], Optional[str]]:
    """
    將 PDF 每頁轉為圖片，嵌入 PPT 幻燈片（等比例縮放填滿）。
    Returns: (pptx_bytes, error_message)
    """
    _safe_imports()
    if not _pptx:
        return None, "未安裝 python-pptx，請執行：pip install python-pptx"
    if not _pdf2image:
        return None, "未安裝 pdf2image（需同時安裝 poppler）。請執行：pip install pdf2image"
    if not _pil:
        return None, "未安裝 Pillow"

    try:
        Presentation, Inches, Pt = _pptx
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        images = _pdf2image(pdf_bytes, dpi=150)
        total = len(images)
        slide_w_inch = 10.0
        slide_h_inch = 7.5
        for i, img in enumerate(images):
            if progress_callback:
                progress_callback((i + 1) / total)

            slide_layout = prs.slide_layouts[6]  # 空白
            slide = prs.slides.add_slide(slide_layout)
            img_w, img_h = img.size
            # 等比例縮放以填滿頁面（使用 max scale）
            img_w_inch = img_w / 96.0
            img_h_inch = img_h / 96.0
            scale_w = slide_w_inch / img_w_inch
            scale_h = slide_h_inch / img_h_inch
            scale = max(scale_w, scale_h)
            new_w_inch = img_w_inch * scale
            new_h_inch = img_h_inch * scale
            left_inch = (slide_w_inch - new_w_inch) / 2
            top_inch = (slide_h_inch - new_h_inch) / 2
            buf_img = io.BytesIO()
            img.save(buf_img, format="PNG")
            buf_img.seek(0)
            slide.shapes.add_picture(buf_img, Inches(left_inch), Inches(top_inch), width=Inches(new_w_inch), height=Inches(new_h_inch))

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read(), None
    except Exception as e:
        err_msg = str(e)
        if "encrypted" in err_msg.lower() or "password" in err_msg.lower():
            return None, "PDF 已加密或受密碼保護，無法讀取"
        if "poppler" in err_msg.lower() or "convert" in err_msg.lower():
            return None, "pdf2image 需要 poppler。Windows 可安裝 poppler 或使用 conda：conda install -c conda-forge poppler"
        return None, f"轉換失敗：{err_msg}"


def pdf_to_images(
    pdf_bytes: bytes,
    fmt: str = "png",
    dpi: int = 200,
    progress_callback=None
) -> Tuple[Optional[bytes], Optional[bytes], Optional[str]]:
    """
    將 PDF 每頁轉為圖片，壓縮為 ZIP。
    Returns: (zip_bytes, first_image_bytes_for_preview, error_message)
    """
    _safe_imports()
    if not _pdf2image:
        return None, None, "未安裝 pdf2image（需同時安裝 poppler）。請執行：pip install pdf2image"
    if not _pil:
        return None, None, "未安裝 Pillow"

    try:
        images = _pdf2image(pdf_bytes, dpi=dpi)
        total = len(images)
        first_img_bytes = None
        zip_buf = io.BytesIO()
        ext = "png" if fmt.lower() == "png" else "jpg"
        save_fmt = "PNG" if ext == "png" else "JPEG"

        with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, img in enumerate(images):
                if progress_callback:
                    progress_callback((i + 1) / total)
                img_buf = io.BytesIO()
                if ext == "jpg":
                    if img.mode in ("RGBA", "P"):
                        img = img.convert("RGB")
                    img.save(img_buf, format=save_fmt, quality=90)
                else:
                    img.save(img_buf, format=save_fmt)
                img_buf.seek(0)
                data = img_buf.read()
                zf.writestr(f"page_{i+1:04d}.{ext}", data)
                if i == 0:
                    first_img_bytes = data

        zip_buf.seek(0)
        return zip_buf.read(), first_img_bytes, None
    except Exception as e:
        err_msg = str(e)
        if "encrypted" in err_msg.lower() or "password" in err_msg.lower():
            return None, None, "PDF 已加密或受密碼保護，無法讀取"
        if "poppler" in err_msg.lower():
            return None, None, "pdf2image 需要 poppler。Windows 可安裝 poppler 或使用 conda"
        return None, None, f"轉換失敗：{err_msg}"


def _docx_text_length(docx_bytes: bytes) -> int:
    """檢查 Word 檔內文字總長度，用於判斷是否為空檔（掃描檔轉換常無內容）。"""
    try:
        from docx import Document
        doc = Document(io.BytesIO(docx_bytes))
        return sum(len(p.text) for p in doc.paragraphs) + sum(
            sum(len(c.text) for c in row.cells) for table in doc.tables for row in table.rows
        )
    except Exception:
        return 0


def pdf_to_word_with_tesseract(
    pdf_bytes: bytes,
    lang: str = "chi_tra+eng",
    dpi: int = 200,
    progress_callback=None,
) -> Tuple[Optional[bytes], Optional[str]]:
    """
    使用 Tesseract OCR 將掃描檔 PDF 轉為可編輯 Word。
    僅產出純文字（無樣式、無獨立圖片）。若要「樣式／字體／圖片皆可編輯」，請用一般模式（pdf2docx）處理文字型 PDF。
    Returns: (docx_bytes, error_message)
    """
    _safe_imports()
    if not _pdf2image:
        return None, "未安裝 pdf2image（需 poppler）。請執行：pip install pdf2image"
    if not _pil:
        return None, "未安裝 Pillow"
    if not _python_docx:
        return None, "未安裝 python-docx，請執行：pip install python-docx"

    try:
        import pytesseract
    except ImportError:
        return None, "未安裝 pytesseract。若使用 venv，請先 source venv/bin/activate 再 pip install pytesseract；或執行 venv/bin/pip install pytesseract"

    try:
        images = _pdf2image(pdf_bytes, dpi=dpi)
        total = len(images)
        if total == 0:
            return None, "PDF 中無頁面"

        Document, Pt = _python_docx
        doc = Document()

        for i, img in enumerate(images):
            if progress_callback:
                progress_callback((i + 1) / total)
            if img.mode != "RGB":
                img = img.convert("RGB")
            text = pytesseract.image_to_string(img, lang=lang)
            for para in (text or "").split("\n\n"):
                para = para.strip()
                if para:
                    p = doc.add_paragraph(para)
                    p.paragraph_format.space_after = Pt(6)
            if i < total - 1:
                doc.add_page_break()

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.read(), None
    except Exception as e:
        err_msg = str(e)
        if "tesseract" in err_msg.lower() or "is not installed" in err_msg.lower():
            return None, "請安裝 Tesseract OCR：Windows 從 https://github.com/UB-Mannheim/tesseract/wiki 下載；Ubuntu 執行 sudo apt install tesseract-ocr tesseract-ocr-chi-tra"
        if "encrypted" in err_msg.lower() or "password" in err_msg.lower():
            return None, "PDF 已加密或受密碼保護，無法讀取"
        if "poppler" in err_msg.lower():
            return None, "pdf2image 需要 poppler。Ubuntu 執行：sudo apt install poppler-utils"
        return None, f"OCR 轉換失敗：{err_msg}"


def pdf_to_word(pdf_bytes: bytes, progress_callback=None) -> Tuple[Optional[bytes], Optional[str]]:
    """
    使用 pdf2docx 將 PDF 轉為 Word。
    Returns: (docx_bytes, error_message)
    """
    _safe_imports()
    if not _pdf2docx:
        return None, "未安裝 pdf2docx，請執行：pip install pdf2docx"

    try:
        Converter = _pdf2docx
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as pdf_tmp:
            pdf_tmp.write(pdf_bytes)
            pdf_path = pdf_tmp.name
        try:
            with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as docx_tmp:
                docx_path = docx_tmp.name
            try:
                cv = Converter(pdf_path, docx_path)
                cv.convert()
                cv.close()
                with open(docx_path, "rb") as f:
                    return f.read(), None
            finally:
                try:
                    os.remove(docx_path)
                except Exception:
                    pass
        finally:
            try:
                os.remove(pdf_path)
            except Exception:
                pass
    except Exception as e:
        err_msg = str(e)
        if "encrypted" in err_msg.lower() or "password" in err_msg.lower():
            return None, "PDF 已加密或受密碼保護，無法讀取"
        return None, f"轉換失敗：{err_msg}"


def pdf_to_word_with_ai_ocr(
    pdf_bytes: bytes,
    api_key: str,
    model_name: str = "gemini-2.0-flash",
    progress_callback=None,
) -> Tuple[Optional[bytes], Optional[str]]:
    """
    使用 Gemini AI Vision 對 PDF 每頁進行 OCR，產出 Word 檔。
    適用掃描檔、圖片型 PDF。
    Returns: (docx_bytes, error_message)
    """
    _safe_imports()
    if not _pdf2image:
        return None, "未安裝 pdf2image（需 poppler）。請執行：pip install pdf2image"
    if not _pil:
        return None, "未安裝 Pillow"
    if not _python_docx:
        return None, "未安裝 python-docx，請執行：pip install python-docx"
    if not _requests:
        return None, "未安裝 requests"
    if not api_key or not api_key.strip():
        return None, "未提供 Gemini API 金鑰"

    try:
        import requests as req
        images = _pdf2image(pdf_bytes, dpi=200)
        total = len(images)
        if total == 0:
            return None, "PDF 中無頁面"

        Document, Pt = _python_docx
        doc = Document()
        prompt = """請將此圖片中的所有文字完整辨識並輸出。
要求：
1. 逐行、逐段輸出，保持原有閱讀順序
2. 保留段落換行（用空行分隔段落）
3. 若為表格，請以空格或 Tab 對齊呈現
4. 純文字輸出，不要加標題或說明
5. 使用繁體中文（若為其他語言則原文輸出）"""

        m_name = model_name if "models/" in model_name else f"models/{model_name}"
        url = f"https://generativelanguage.googleapis.com/v1beta/{m_name}:generateContent?key={api_key.strip()}"

        resample = _pil_lanczos()
        for i, img in enumerate(images):
            if progress_callback:
                progress_callback((i + 1) / total)

            if img.mode != "RGB":
                img = img.convert("RGB")
            if resample is not None:
                img.thumbnail((1920, 1920), resample)
            else:
                img.thumbnail((1920, 1920))
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=90)
            img_b64 = base64.b64encode(buf.getvalue()).decode()

            payload = {
                "contents": [{
                    "parts": [
                        {"text": prompt},
                        {"inlineData": {"mimeType": "image/jpeg", "data": img_b64}},
                    ]
                }],
                "generationConfig": {"temperature": 0.1, "maxOutputTokens": 8192},
            }

            resp = req.post(url, json=payload, timeout=60)
            if resp.status_code != 200:
                return None, f"Gemini API 錯誤: {resp.status_code} {resp.text[:200]}"

            data = resp.json()
            if not data.get("candidates") or not data["candidates"][0].get("content", {}).get("parts"):
                return None, f"第 {i+1} 頁 OCR 無回傳內容"

            text = data["candidates"][0]["content"]["parts"][0].get("text", "").strip()
            if text:
                for para in text.split("\n\n"):
                    para = para.strip()
                    if para:
                        p = doc.add_paragraph(para)
                        p.paragraph_format.space_after = Pt(6)

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)
        return buf.read(), None
    except Exception as e:
        err_msg = str(e)
        if "encrypted" in err_msg.lower() or "password" in err_msg.lower():
            return None, "PDF 已加密或受密碼保護，無法讀取"
        if "poppler" in err_msg.lower():
            return None, "pdf2image 需要 poppler"
        return None, f"AI OCR 轉換失敗：{err_msg}"


# --- AI 高品質版面還原（Gemini Vision 結構化 JSON → Word）---
_AI_LAYOUT_PROMPT = """请将这张 PDF 页面的内容解析成结构化 JSON。只输出一个 JSON 对象，不要其他说明或 markdown 标记。

格式要求：
{
  "page": 1,
  "blocks": [
    {"type": "title", "text": "主标题文字"},
    {"type": "subtitle", "text": "副标题"},
    {"type": "heading", "level": 1, "text": "标题文字"},
    {"type": "paragraph", "text": "段落内容"},
    {"type": "bullet_list", "items": ["项目1", "项目2"]},
    {"type": "table", "header": ["列1", "列2"], "rows": [["a","b"], ["c","d"]]},
    {"type": "image", "id": "img1"}
  ]
}

规则：
1. 按阅读顺序列出所有区块。
2. type 只能是：title, subtitle, heading, paragraph, bullet_list, table, image 之一。
3. 表格用 header（可选）和 rows 二维数组。
4. 图片区块用 type:"image" 和 id:"img1", "img2"... 按出现顺序编号。
5. 只输出上述 JSON，不要 ``` 包裹。"""


def _extract_images_from_pdf_path(pdf_path: str) -> List[List[bytes]]:
    """每頁抽取的圖片 bytes。images_per_page[i] = 第 i 頁的 [img1_bytes, ...]。"""
    _safe_imports()
    if not _pymupdf:
        return []
    out = []
    try:
        doc = _pymupdf.open(pdf_path)
        for page in doc:
            img_list = []
            for img in page.get_images():
                xref = img[0]
                base = doc.extract_image(xref)
                img_list.append(base["image"])
            out.append(img_list)
        doc.close()
    except Exception:
        pass
    return out


def _parse_ai_layout_json(raw: str) -> dict:
    """從模型輸出解析 JSON（允許 ```json ... ``` 包裝）。"""
    raw = raw.strip()
    m = re.search(r"```(?:json)?\s*([\s\S]*?)```", raw)
    if m:
        raw = m.group(1).strip()
    return json.loads(raw)


def _build_docx_from_ai_layout_pages(
    pages_data: List[dict],
    images_per_page: List[List[bytes]],
    page_pil_images: list,
) -> bytes:
    """依每頁 JSON 與圖片建出 .docx bytes。"""
    Document, Pt = _python_docx
    from docx.shared import Inches
    doc = Document()
    for page_idx, page in enumerate(pages_data):
        blocks = page.get("blocks") or []
        page_images = images_per_page[page_idx] if page_idx < len(images_per_page) else []
        pil_page = page_pil_images[page_idx] if page_idx < len(page_pil_images) else None
        img_index = 0
        for block in blocks:
            t = (block.get("type") or "").strip().lower()
            text = (block.get("text") or "").strip()
            if t == "title":
                p = doc.add_paragraph(text)
                p.style = "Title"
            elif t == "subtitle":
                p = doc.add_paragraph(text)
                p.style = "Subtitle"
            elif t == "heading":
                level = block.get("level", 1)
                p = doc.add_paragraph(text)
                p.style = "Heading %d" % min(level, 3)
            elif t == "paragraph":
                p = doc.add_paragraph(text)
                p.paragraph_format.space_after = Pt(6)
            elif t == "bullet_list":
                for item in block.get("items") or []:
                    p = doc.add_paragraph(str(item).strip(), style="List Bullet")
                    p.paragraph_format.space_after = Pt(3)
            elif t == "table":
                header = block.get("header") or []
                rows = block.get("rows") or []
                if header or rows:
                    col_count = max(len(header), max([len(r) for r in rows], default=0)) or 1
                    row_count = (1 if header else 0) + len(rows)
                    table = doc.add_table(rows=row_count, cols=col_count)
                    table.style = "Table Grid"
                    r = 0
                    if header:
                        for c, cell_text in enumerate(header):
                            if c < col_count:
                                table.rows[r].cells[c].text = str(cell_text)
                        r += 1
                    for row in rows:
                        for c, cell_text in enumerate(row):
                            if c < col_count:
                                table.rows[r].cells[c].text = str(cell_text)
                        r += 1
            elif t == "image":
                if img_index < len(page_images):
                    try:
                        doc.add_picture(io.BytesIO(page_images[img_index]), width=Inches(3.0))
                    except Exception:
                        if pil_page:
                            buf = io.BytesIO()
                            pil_page.save(buf, format="PNG")
                            buf.seek(0)
                            doc.add_picture(buf, width=Inches(4.0))
                    img_index += 1
                elif pil_page:
                    buf = io.BytesIO()
                    pil_page.save(buf, format="PNG")
                    buf.seek(0)
                    doc.add_picture(buf, width=Inches(4.0))
                    img_index += 1
        if page_idx < len(pages_data) - 1:
            doc.add_page_break()
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def pdf_to_word_with_ai_layout(
    pdf_bytes: bytes,
    api_key: str,
    model_name: str = "gemini-2.0-flash",
    progress_callback=None,
) -> Tuple[Optional[bytes], Optional[str]]:
    """
    使用 Gemini Vision 解析 PDF 每頁結構（標題、段落、清單、表格、圖片），
    以 python-docx 重建可編輯 Word（樣式、字體、圖片皆可編輯）。適用掃描檔或文字型 PDF。
    Returns: (docx_bytes, error_message)
    """
    _safe_imports()
    if not _pdf2image:
        return None, "未安裝 pdf2image（需 poppler）。請執行：pip install pdf2image"
    if not _pil:
        return None, "未安裝 Pillow"
    if not _python_docx:
        return None, "未安裝 python-docx，請執行：pip install python-docx"
    if not _requests:
        return None, "未安裝 requests"
    if not api_key or not api_key.strip():
        return None, "未提供 Gemini API 金鑰"

    try:
        req = _requests
        images = _pdf2image(pdf_bytes, dpi=200)
        total = len(images)
        if total == 0:
            return None, "PDF 中無頁面"

        pdf_path = None
        try:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tf:
                tf.write(pdf_bytes)
                pdf_path = tf.name
            images_per_page = _extract_images_from_pdf_path(pdf_path)
        finally:
            if pdf_path and os.path.isfile(pdf_path):
                try:
                    os.remove(pdf_path)
                except Exception:
                    pass

        m_name = model_name if "models/" in model_name else ("models/%s" % model_name)
        url = "https://generativelanguage.googleapis.com/v1beta/%s:generateContent?key=%s" % (m_name, api_key.strip())
        pages_data = []
        resample = _pil_lanczos()
        for i, img in enumerate(images):
            if progress_callback:
                progress_callback((i + 1) / total)
            if img.mode != "RGB":
                img = img.convert("RGB")
            if resample is not None:
                img.thumbnail((1600, 1600), resample)
            else:
                img.thumbnail((1600, 1600))
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=88)
            img_b64 = base64.b64encode(buf.getvalue()).decode()
            payload = {
                "contents": [{
                    "parts": [
                        {"text": _AI_LAYOUT_PROMPT},
                        {"inlineData": {"mimeType": "image/jpeg", "data": img_b64}},
                    ]
                }],
                "generationConfig": {"temperature": 0.1, "maxOutputTokens": 8192, "responseMimeType": "application/json"},
            }
            resp = req.post(url, json=payload, timeout=120)
            if resp.status_code != 200:
                return None, "Gemini API 錯誤: %s %s" % (resp.status_code, resp.text[:200])
            data = resp.json()
            if not data.get("candidates") or not data["candidates"][0].get("content", {}).get("parts"):
                return None, "第 %d 頁 AI 未回傳內容" % (i + 1)
            raw = data["candidates"][0]["content"]["parts"][0].get("text", "").strip()
            try:
                page_json = _parse_ai_layout_json(raw)
            except json.JSONDecodeError:
                page_json = {"page": i + 1, "blocks": [{"type": "paragraph", "text": raw[:5000]}]}
            pages_data.append(page_json)

        docx_bytes = _build_docx_from_ai_layout_pages(pages_data, images_per_page, images)
        return docx_bytes, None
    except Exception as e:
        err_msg = str(e)
        if "encrypted" in err_msg.lower() or "password" in err_msg.lower():
            return None, "PDF 已加密或受密碼保護，無法讀取"
        if "poppler" in err_msg.lower():
            return None, "pdf2image 需要 poppler"
        return None, "AI 高品質排版轉換失敗：%s" % err_msg
